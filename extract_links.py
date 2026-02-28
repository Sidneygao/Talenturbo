from pptx import Presentation
import sys

pptx_file = sys.argv[1] if len(sys.argv) > 1 else "Modular-HR-Digital-Platform-Centered-on-HR-CORE.pptx"

try:
    prs = Presentation(pptx_file)
    print(f"=== Links from {pptx_file} ===\n")
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_links = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Check for hyperlinks in text
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if hasattr(run, "hyperlink") and run.hyperlink.address:
                                slide_links.append({
                                    'text': run.text,
                                    'url': run.hyperlink.address,
                                    'context': shape.text[:200]
                                })
        
        if slide_links:
            print(f"Slide {slide_num}:")
            for link in slide_links:
                print(f"  Text: {link['text']}")
                print(f"  URL: {link['url']}")
                print(f"  Context: {link['context']}")
                print()

except Exception as e:
    print(f"Error: {e}")
