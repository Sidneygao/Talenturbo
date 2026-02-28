#!/usr/bin/env python
# -*- coding: utf-8 -*-
from pptx import Presentation
import sys

def extract_hyperlinks(pptx_path):
    prs = Presentation(pptx_path)
    
    print(f"\n{'='*80}")
    print(f"Extracting hyperlinks from: {pptx_path}")
    print(f"{'='*80}\n")
    
    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_title = ""
        
        # Try to get slide title
        if slide.shapes.title:
            slide_title = slide.shapes.title.text.strip()
        
        links_found = []
        
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
                
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if hasattr(run, 'hyperlink') and run.hyperlink and run.hyperlink.address:
                        link_text = run.text.strip()
                        link_url = run.hyperlink.address
                        
                        # Get surrounding context
                        full_text = shape.text.strip()
                        
                        links_found.append({
                            'text': link_text,
                            'url': link_url,
                            'context': full_text[:300]
                        })
        
        if links_found:
            print(f"Slide {slide_idx}: {slide_title if slide_title else '(No title)'}")
            print("-" * 80)
            for link in links_found:
                print(f"  Link Text: {link['text']}")
                print(f"  URL: {link['url']}")
                print(f"  Context: {link['context']}")
                print()

if __name__ == "__main__":
    pptx_file = sys.argv[1] if len(sys.argv) > 1 else "Modular-HR-Digital-Platform-Centered-on-HR-CORE.pptx"
    
    try:
        extract_hyperlinks(pptx_file)
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
