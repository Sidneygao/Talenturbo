import json
import os
import sys

from pptx import Presentation


def iter_shape_text_with_links(shape):
    items = []

    # Shape-level hyperlink (click action)
    try:
        addr = shape.click_action.hyperlink.address
        if addr:
            items.append({"type": "shape", "text": getattr(shape, "text", "") or "", "url": addr})
    except Exception:
        pass

    if not getattr(shape, "has_text_frame", False):
        return items

    tf = shape.text_frame
    for p in tf.paragraphs:
        runs = getattr(p, "runs", None)
        if not runs:
            t = p.text.strip()
            if t:
                items.append({"type": "paragraph", "text": t, "url": None})
            continue

        for r in runs:
            t = (r.text or "").strip()
            if not t:
                continue
            url = None
            try:
                url = r.hyperlink.address
            except Exception:
                url = None
            items.append({"type": "run", "text": t, "url": url})

    return items


def extract(pptx_path: str):
    prs = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        shapes = []
        all_text = []
        links = []

        for s in slide.shapes:
            shape_entry = {
                "name": getattr(s, "name", None),
                "shape_type": str(getattr(s, "shape_type", None)),
                "text": (getattr(s, "text", "") or "").strip(),
                "text_items": [],
            }

            text_items = iter_shape_text_with_links(s)
            shape_entry["text_items"] = text_items
            shapes.append(shape_entry)

            for it in text_items:
                if it.get("text"):
                    all_text.append(it["text"])
                if it.get("url"):
                    links.append({"text": it.get("text") or "", "url": it["url"]})

        # Heuristic title: first non-empty text shape
        title = ""
        for sh in shapes:
            t = (sh.get("text") or "").strip()
            if t:
                title = t.splitlines()[0].strip()
                break

        slides.append(
            {
                "index": i,
                "title": title,
                "texts": [t for t in all_text if t],
                "links": links,
                "shapes": shapes,
            }
        )

    return {"pptx": os.path.basename(pptx_path), "slide_count": len(slides), "slides": slides}


def main():
    if len(sys.argv) < 2:
        print("Usage: py extract_pptx.py <path-to-pptx>", file=sys.stderr)
        raise SystemExit(2)

    pptx_path = sys.argv[1]
    if not os.path.exists(pptx_path):
        print(f"File not found: {pptx_path}", file=sys.stderr)
        raise SystemExit(1)

    data = extract(pptx_path)
    print(json.dumps(data, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()

